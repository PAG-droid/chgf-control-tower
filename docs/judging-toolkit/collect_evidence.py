#!/usr/bin/env python
"""Deterministic evidence collection for hackathon rubric scoring.

Reads a submission's deck (.pptx) and one or more local git repos, and emits a single
JSON evidence file. Scores NOTHING and judges NOTHING -- it only produces citable facts
(slide numbers, commit SHAs, file paths) so that every band assigned later can point at
one. Anything this script cannot observe is emitted as null/empty rather than guessed.

Usage (one segment, no shell metacharacters):
  python collect_evidence.py --project "NAME" --deck deck.pptx --repo C:/path/to/repo
      --since 2026-08-11 --until 2026-08-15 --out evidence.json

  --repo may be repeated. Remote repos must be cloned locally first; this script makes
  no network calls.
"""
from __future__ import annotations

import argparse
import json
import os
import re
import subprocess
import sys
from datetime import datetime, timezone

# ---------------------------------------------------------------- deck ------------

URL_RE = re.compile(r"https?://[^\s)>\]\"']+")
# Quantified claims: money, percentages, multipliers, counts with a unit noun.
CLAIM_RE = re.compile(
    r"(\$\s?[\d,.]+\s?(?:[kKmMbB]|billion|million|thousand)?"
    r"|[\d,.]+\s?%"
    r"|[\d,.]+\s?[xX]\b"
    r"|\b[\d,.]+\s+(?:users?|people|hours?|days?|weeks?|FTEs?|teams?|countries|"
    r"investments?|grants?|documents?|minutes?|seconds?|queries|records?|rows?)\b)",
    re.IGNORECASE,
)
FORWARD_LOOK_RE = re.compile(
    r"\b(next steps?|roadmap|future|phase 2|what'?s next|coming soon|planned|vision|"
    r"if we had more time|backlog)\b",
    re.IGNORECASE,
)
DEMO_RE = re.compile(r"\b(demo|screenshot|live|walkthrough|video|try it|deployed)\b", re.IGNORECASE)

# ------------------------------------------------- superlative-award counters ----
# All of the below feed the FUN awards only. They never touch a 1-3 rubric band.
# Lexicons are deliberately short, visible and tunable -- edit them for your crowd.

# NOTE: arrow blocks (U+2190-21FF, U+2900-297F) are deliberately EXCLUDED. A first cut
# included them and scored 3,287 "emoji" on a docs-heavy repo -- every "was -> now" arrow
# and box-drawing rule counted. Typographic marks that render as text, not emoji, are
# stripped by TYPOGRAPHY_DENY below for the same reason.
EMOJI_RE = re.compile(
    "[" "\U0001F000-\U0001FAFF" "\U0001F1E6-\U0001F1FF" "☀-⛿" "✀-➿" "⬀-⯿" "]"
)
TYPOGRAPHY_DENY = set("★☆▪▫●○■□◆◇▲▼◄►✓✔✗✘✂✎✏✁⬛⬜⬝⬞⯀⯁")
# ALL-CAPS tokens of 2-6 letters: acronym density proxy.
ACRONYM_RE = re.compile(r"\b[A-Z]{2,6}\b")

GEEK_LEXICON = [
    "regex", "ast", "monad", "functor", "idempotent", "dag", "graph theory", "eigen",
    "bayes", "markov", "big-o", "o(n", "assembly", "bitmask", "compiler", "parser",
    "dsl", "wasm", "cuda", "tensor", "embedding", "vector space", "cypher", "simd",
    "mutex", "closure", "recursion", "lambda calculus", "category theory", "np-hard",
    "hash", "entropy", "gradient", "jacobian", "transformer", "attention head",
    "quantiz", "tokenizer", "kernel", "bare metal", "protobuf", "grpc", "systemd",
    "kubernetes", "bitwise", "traversal", "topolog", "heuristic", "combinatorial",
    "deterministic", "async", "concurrency", "state machine", "bloom filter",
]
NERD_LEXICON = [
    "star wars", "jedi", "sith", "yoda", "lord of the rings", "frodo", "gandalf",
    "mordor", "hobbit", "sauron", "one ring", "pokemon", "pikachu", "star trek",
    "klingon", "tribble", "borg", "tardis", "doctor who", "hitchhiker", "don't panic",
    "answer to life", "dungeons", "d&d", "mana", "xp", "level up", "wizard", "dragon",
    "kraken", "phoenix", "ninja", "sensei", "the matrix", "red pill", "skynet",
    "hal 9000", "xkcd", "thanos", "batman", "tetris", "mario", "zelda", "minecraft",
    "discworld", "dune", "melange", "cthulhu", "hogwarts", "muggle", "may the force",
    "easter egg", "konami", "rubber duck", "yak shave", "bikeshed",
]
# Plain-business register. "Normalist" = high here, low on the two above.
NORMIE_LEXICON = [
    "stakeholder", "roi", "alignment", "workflow", "process", "team", "timeline",
    "budget", "adoption", "training", "rollout", "pilot", "hours saved", "manual",
    "spreadsheet", "email", "meeting", "approval", "review", "deadline", "colleague",
    "customer", "handoff", "onboarding", "reporting", "dashboard", "efficiency",
    "turnaround", "backlog of work", "sign-off", "best practice", "value",
]
TOOLING_HINT_RE = re.compile(
    r"^(\.claude/|\.github/|docs?/|\.vscode/|\.devcontainer/)|"
    r"(^|/)(Dockerfile|docker-compose\.ya?ml|Makefile|\.env\.(example|sample)|"
    r"CLAUDE\.md|AGENTS\.md|README\.md|CONTRIBUTING\.md|\.gitignore|"
    r"requirements\.txt|pyproject\.toml|package\.json|settings\.json|\.mcp\.json)$|"
    r"\.(ya?ml|toml|ini|cfg)$",
)
TEXT_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".md", ".txt", ".ps1", ".sql", ".java", ".go",
    ".rb", ".sh", ".yml", ".yaml", ".json", ".html", ".css", ".cs", ".cpp", ".c", ".h",
    ".rs", ".php", ".r", ".ipynb", ".toml",
}
TODO_RE = re.compile(r"\b(TODO|FIXME|HACK|XXX)\b")
SHORTSTAT_RE = re.compile(r"(\d+) files? changed(?:, (\d+) insertions?\(\+\))?(?:, (\d+) deletions?\(-\))?")


def _lex_hits(text, lexicon):
    """Return {term: count} for lexicon terms present in text (lowercased substring)."""
    low = text.lower()
    return {t: low.count(t) for t in lexicon if t in low}


def _emoji(text):
    found = [c for c in EMOJI_RE.findall(text) if c not in TYPOGRAPHY_DENY]
    tally = {}
    for e in found:
        tally[e] = tally.get(e, 0) + 1
    return len(found), tally


def _shape_text(shape, out):
    """Collect text from a shape, recursing into groups and tables."""
    try:
        if shape.shape_type is not None and int(shape.shape_type) == 6:  # GROUP
            for sub in shape.shapes:
                _shape_text(sub, out)
            return
    except Exception:
        pass
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    out.append(cell.text.strip())
        return
    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
        out.append(shape.text_frame.text.strip())


def _shape_links(shape, out):
    try:
        if shape.shape_type is not None and int(shape.shape_type) == 6:
            for sub in shape.shapes:
                _shape_links(sub, out)
            return
    except Exception:
        pass
    try:
        addr = shape.click_action.hyperlink.address
        if addr:
            out.append(addr)
    except Exception:
        pass
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.hyperlink and run.hyperlink.address:
                        out.append(run.hyperlink.address)
                except Exception:
                    pass


def read_deck(path):
    if not path:
        return None
    if not os.path.isfile(path):
        return {"error": f"deck not found: {path}"}
    if not path.lower().endswith(".pptx"):
        return {"error": f"not a .pptx ({os.path.splitext(path)[1]}) -- read it with the Read tool instead"}
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed"}

    prs = Presentation(path)
    slides, all_links, claims, media_total, picture_total = [], [], [], 0, 0
    for idx, slide in enumerate(prs.slides, start=1):
        texts, links = [], []
        pictures = charts = tables = media = 0
        for shape in slide.shapes:
            _shape_text(shape, texts)
            _shape_links(shape, links)
            try:
                st = int(shape.shape_type) if shape.shape_type is not None else -1
            except Exception:
                st = -1
            if st == 13:
                pictures += 1
            elif st == 16:
                media += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1
        body = "\n".join(texts)
        notes = ""
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        links += URL_RE.findall(body + "\n" + notes)
        slide_claims = [c.strip() for c in CLAIM_RE.findall(body)]
        slide_emoji, _ = _emoji(body)
        slides.append(
            {
                "slide": idx,
                "title": texts[0][:120] if texts else "",
                "text": body,
                "notes": notes,
                "word_count": len(body.split()),
                "pictures": pictures,
                "charts": charts,
                "tables": tables,
                "media": media,
                "links": sorted(set(links)),
                "quantified_claims": slide_claims,
                "emoji_count": slide_emoji,
                "forward_looking": bool(FORWARD_LOOK_RE.search(body)),
                "demo_signal": bool(DEMO_RE.search(body + " " + notes)),
            }
        )
        all_links += links
        claims += [{"slide": idx, "claim": c} for c in slide_claims]
        picture_total += pictures
        media_total += media

    deck_text = "\n".join(s["text"] for s in slides)
    notes_text = "\n".join(s["notes"] for s in slides)
    emoji_total, emoji_tally = _emoji(deck_text + "\n" + notes_text)
    body_words = sum(s["word_count"] for s in slides)
    notes_words = len(notes_text.split())
    acronyms = ACRONYM_RE.findall(deck_text)
    geek = _lex_hits(deck_text, GEEK_LEXICON)
    nerd = _lex_hits(deck_text, NERD_LEXICON)
    normie = _lex_hits(deck_text, NORMIE_LEXICON)

    return {
        "path": os.path.abspath(path),
        "slide_count": len(slides),
        "total_words": body_words,
        "notes_words": notes_words,
        "picture_count": picture_total,
        "embedded_media_count": media_total,
        "slides_with_notes": sum(1 for s in slides if s["notes"]),
        "unique_links": sorted(set(all_links)),
        "quantified_claims": claims,
        "forward_looking_slides": [s["slide"] for s in slides if s["forward_looking"]],
        "demo_signal_slides": [s["slide"] for s in slides if s["demo_signal"]],
        # --- superlative-award counters (fun awards only) ---
        "emoji_count": emoji_total,
        "emoji_tally": dict(sorted(emoji_tally.items(), key=lambda kv: -kv[1])[:20]),
        "emoji_per_slide": round(emoji_total / len(slides), 2) if slides else 0,
        "est_speaking_minutes": round(body_words / 130.0, 1),
        "est_speaking_minutes_with_notes": round((body_words + notes_words) / 130.0, 1),
        "words_per_slide": round(body_words / len(slides), 1) if slides else 0,
        "densest_slide": max(slides, key=lambda s: s["word_count"])["slide"] if slides else None,
        "densest_slide_words": max((s["word_count"] for s in slides), default=0),
        "acronym_tokens": len(acronyms),
        "unique_acronyms": sorted(set(acronyms))[:40],
        "geek_lexicon_hits": geek,
        "geek_score": sum(geek.values()),
        "nerd_lexicon_hits": nerd,
        "nerd_score": sum(nerd.values()),
        "normie_lexicon_hits": normie,
        "normie_score": sum(normie.values()),
        "slides": slides,
    }


# ---------------------------------------------------------------- repo ------------

CLAUDE_ARTIFACTS = {
    "claude_md": ["CLAUDE.md", ".claude/CLAUDE.md"],
    "settings": [".claude/settings.json", ".claude/settings.local.json"],
    "mcp_config": [".mcp.json", ".claude/mcp.json"],
    "agents_dir": [".claude/agents"],
    "skills_dir": [".claude/skills", "skills"],
    "commands_dir": [".claude/commands", "commands"],
    "hooks_dir": [".claude/hooks"],
    "plugin_manifest": [".claude-plugin/plugin.json", "plugin.json"],
}
DURABILITY_ARTIFACTS = {
    "readme": ["README.md", "README.rst", "readme.md"],
    "license": ["LICENSE", "LICENSE.md", "LICENCE"],
    "ci": [".github/workflows", ".gitlab-ci.yml", "azure-pipelines.yml"],
    "container": ["Dockerfile", "docker-compose.yml", "compose.yaml"],
    "deps": ["requirements.txt", "pyproject.toml", "package.json", "Gemfile", "go.mod", "pom.xml"],
    "config_example": [".env.example", ".env.sample", "config.example.json"],
    "docs_dir": ["docs"],
    "install_docs": ["INSTALL.md", "SETUP.md", "docs/ONBOARDING.md", "CONTRIBUTING.md"],
    # Found missing during a real dry run: a submission shipped .githooks/pre-commit and
    # .githooks/pre-push as its validation gate and neither table saw them, so genuine
    # Staying Power evidence had to be cited by hand.
    "git_hooks": [".githooks", ".pre-commit-config.yaml", ".husky"],
}
TEST_HINT_RE = re.compile(r"(^|/)(tests?|spec|__tests__)(/|$)|(^|/)test_[^/]+\.py$|[^/]+_test\.(py|go|js|ts)$|[^/]+\.(test|spec)\.(js|ts|tsx|jsx)$")
CLAUDE_TRAILER_RE = re.compile(r"co-authored-by:\s*claude|claude\.com/claude-code|generated with \[claude", re.IGNORECASE)
SKIP_DIRS = {".git", "node_modules", ".venv", "venv", "__pycache__", "dist", "build", ".next", "target"}
US = "\x1f"
RS = "\x1e"


def _git(repo, args):
    try:
        p = subprocess.run(["git"] + args, cwd=repo, capture_output=True, text=True,
                           encoding="utf-8", errors="replace", timeout=180)
    except (OSError, subprocess.TimeoutExpired) as exc:
        return None, str(exc)
    if p.returncode != 0:
        return None, (p.stderr or "").strip()
    return p.stdout, None


def _ancestors(path):
    parts = path.split("/")
    return {"/".join(parts[:i]) for i in range(1, len(parts))}


def _read_blob(repo, rel, ref_tip):
    """File content from the ref (git show) or the working tree. Returns None if absent."""
    if ref_tip:
        out, _ = _git(repo, ["show", f"{ref_tip}:{rel}"])
        return out
    full = os.path.join(repo, rel.replace("/", os.sep))
    try:
        if os.path.getsize(full) > 512 * 1024:
            return None
        with open(full, "r", encoding="utf-8", errors="replace") as fh:
            return fh.read()
    except OSError:
        return None


def _exists_any(repo, candidates, prefixes=None, tracked=None):
    """Find an artifact, preferring the SUBMISSION's own subdir over the repo root.

    On a shared hackathon monorepo a root-level CLAUDE.md belongs to the organisers, not
    to this team -- so the scope is recorded and the caller must not credit a root-scoped
    artifact as the team's own craft.
    """
    scopes = [(p.rstrip("/*"), "submission") for p in (prefixes or [])] + [("", "repo-root")]
    tset = set(tracked or [])
    tdirs = {d for f in tset for d in _ancestors(f)}
    for prefix, scope in scopes:
        for c in candidates:
            rel = f"{prefix}/{c}" if prefix else c
            if tracked is not None:
                # Ref-scoped: the ref's tree is the truth, not whatever is checked out.
                if rel in tset:
                    return {"path": rel, "is_dir": False, "entries": None, "scope": scope}
                if rel in tdirs:
                    n = len({f[len(rel) + 1:].split("/")[0] for f in tset if f.startswith(rel + "/")})
                    return {"path": rel, "is_dir": True, "entries": n, "scope": scope}
                continue
            full = os.path.join(repo, rel.replace("/", os.sep))
            if os.path.exists(full):
                if os.path.isdir(full):
                    try:
                        n = len([e for e in os.listdir(full) if not e.startswith(".")])
                    except OSError:
                        n = 0
                    return {"path": rel, "is_dir": True, "entries": n, "scope": scope}
                return {"path": rel, "is_dir": False, "entries": None, "scope": scope}
    return None


def _tracked_files(repo, paths, ref_tip=None):
    """Tracked files for the submission -- from the REF when one is given.

    ls-files reads the index/checkout, so on a shared monorepo a --ref branch that is not
    checked out silently returned the CHECKOUT's file list: two branch submissions both
    reported "5 test files" because they inherited main's tree. ls-tree reads the ref.
    """
    if ref_tip:
        args = ["ls-tree", "-r", "--name-only", ref_tip] + (["--"] + list(paths) if paths else [])
    else:
        args = ["ls-files"] + (["--"] + list(paths) if paths else [])
    out, err = _git(repo, args)
    if out is None:
        return [], err
    return [f for f in out.splitlines() if f], None


def read_repo(repo, since, until, paths=None, ref=None):
    """Collect one repo's evidence, optionally SCOPED to pathspecs.

    Scoping is not optional hygiene on a shared hackathon monorepo -- without it every
    team inherits every other team's commits, authors and diffstat.
    """
    paths = list(paths or [])
    repo = os.path.abspath(repo)
    # --all spans every branch in the repo. On a shared monorepo where teams submit on
    # their OWN branch, that blends every team's history together -- so a branch-scoped
    # submission MUST pass ref=origin/<branch> instead.
    refs = [ref] if ref else ["--all"]
    # A ref may be a RANGE (origin/main..origin/theirs) to isolate a team's own commits on
    # a shared branch. git log/rev-list take ranges; the file side needs the TIP only.
    ref_tip = ref.split("..")[-1] if ref else None
    result = {"path": repo, "scoped_to": paths or None, "ref": ref or "--all",
              "window": {"since": since, "until": until}}
    if not os.path.isdir(os.path.join(repo, ".git")):
        # could still be a worktree/submodule -- ask git
        out, _ = _git(repo, ["rev-parse", "--is-inside-work-tree"])
        if not out or out.strip() != "true":
            result["error"] = "not a git repository"
            return result

    # Validate the pathspec FIRST. A pathspec matching nothing yields a silently EMPTY
    # submission, which reads as "they built nothing" rather than "wrong path" -- and it
    # happens routinely, because --paths is written for the shared monorepo but applies to
    # every --repo, including a team's standalone repo.
    if paths:
        scoped, _ = _tracked_files(repo, paths, ref_tip)
        if not scoped:
            unscoped, _ = _tracked_files(repo, None, ref_tip)
            if unscoped:
                result["scope_warning"] = (
                    f"pathspec {paths} matched no tracked file here; collected UNSCOPED "
                    f"instead. Check --paths for this repo.")
                paths = []
                result["scoped_to"] = None
            else:
                result["scope_warning"] = "repository has no tracked files at all"

    rng = []
    if since:
        rng.append(f"--since={since}")
    if until:
        rng.append(f"--until={until}")
    ps = (["--"] + paths) if paths else []

    fmt = US.join(["%H", "%an", "%ae", "%aI", "%s", "%b"]) + RS
    out, err = _git(repo, ["log"] + refs + ["--no-merges", f"--pretty=format:{fmt}"] + rng + ps)
    if out is None:
        result["error"] = f"git log failed: {err}"
        return result

    commits, authors, claude_commits = [], {}, 0
    hour_hist = {h: 0 for h in range(24)}
    msg_emoji_total, msg_emoji_tally, emoji_commits = 0, {}, 0
    reverts, longest_msg = 0, {"sha": None, "chars": 0}
    for rec in out.split(RS):
        rec = rec.strip("\n")
        if not rec.strip():
            continue
        parts = rec.split(US)
        if len(parts) < 6:
            continue
        sha, an, ae, aiso, subj, body = parts[0], parts[1], parts[2], parts[3], parts[4], parts[5]
        is_claude = bool(CLAUDE_TRAILER_RE.search(subj + "\n" + body))
        claude_commits += int(is_claude)
        authors[f"{an} <{ae}>"] = authors.get(f"{an} <{ae}>", 0) + 1
        # Author-local hour, straight off the %aI offset -- a 03:00 commit is 03:00 where
        # the author was sitting, which is the whole point of the nocturnal award.
        try:
            hour_hist[int(aiso[11:13])] += 1
        except (ValueError, IndexError, KeyError):
            pass
        n_em, tally = _emoji(subj + "\n" + body)
        msg_emoji_total += n_em
        emoji_commits += int(n_em > 0)
        for k, v in tally.items():
            msg_emoji_tally[k] = msg_emoji_tally.get(k, 0) + v
        if subj.lower().startswith("revert"):
            reverts += 1
        msg_len = len(subj) + len(body)
        if msg_len > longest_msg["chars"]:
            longest_msg = {"sha": sha[:10], "chars": msg_len, "subject": subj[:120]}
        commits.append({"sha": sha[:10], "author": an, "date": aiso,
                        "subject": subj[:160], "claude_trailer": is_claude,
                        "emoji": n_em})

    # --no-merges here MUST match the log above. Counting merges in the total while
    # excluding them from the windowed log made commits_outside_window overstate, and on a
    # merge-heavy repo it can go negative.
    total_all_time, _ = _git(repo, ["rev-list"] + refs + ["--count", "--no-merges"] + ps)
    merges, _ = _git(repo, ["rev-list"] + refs + ["--count", "--merges"] + rng + ps)

    # LOC + files touched inside the window
    ins = dels = 0
    touched = {}
    numstat, _ = _git(repo, ["log"] + refs + ["--no-merges", "--numstat", "--pretty=format:"] + rng + ps)
    for line in (numstat or "").splitlines():
        cols = line.split("\t")
        if len(cols) != 3:
            continue
        a, d, path = cols
        ins += int(a) if a.isdigit() else 0
        dels += int(d) if d.isdigit() else 0
        touched[path] = touched.get(path, 0) + 1

    dates = sorted(c["date"] for c in commits if c["date"])
    active_days = sorted({d[:10] for d in dates})
    top_dirs = {}
    for path in touched:
        top = path.split("/")[0] if "/" in path else "(root)"
        top_dirs[top] = top_dirs.get(top, 0) + 1

    tracked, tracked_err = _tracked_files(repo, paths, ref_tip)
    test_files = [f for f in tracked if TEST_HINT_RE.search(f)]
    ext_counts = {}
    for f in tracked:
        ext = os.path.splitext(f)[1].lower() or "(none)"
        ext_counts[ext] = ext_counts.get(ext, 0) + 1

    # --- superlative counters: biggest single commit, yak-shave ratio, TODOs, deps ---
    biggest = {"sha": None, "files": 0, "churn": 0}
    ss, _ = _git(repo, ["log"] + refs + ["--no-merges", "--shortstat", f"--pretty=format:{RS}%H"] + rng + ps)
    for chunk in (ss or "").split(RS):
        lines = [l for l in chunk.splitlines() if l.strip()]
        if not lines:
            continue
        m = SHORTSTAT_RE.search("\n".join(lines[1:]))
        if not m:
            continue
        files = int(m.group(1) or 0)
        churn = int(m.group(2) or 0) + int(m.group(3) or 0)
        if churn > biggest["churn"]:
            biggest = {"sha": lines[0][:10], "files": files, "churn": churn}

    tooling_touches = sum(n for p, n in touched.items() if TOOLING_HINT_RE.search(p))
    product_touches = sum(touched.values()) - tooling_touches

    todo_count, emoji_in_code, scanned = 0, 0, 0
    for f in tracked:
        if scanned >= 2000:
            break
        if os.path.splitext(f)[1].lower() not in TEXT_EXTS:
            continue
        blob = _read_blob(repo, f, ref_tip)
        if blob is None:
            continue
        scanned += 1
        # TODO markers only in CODE. Counting them in .md turns a docs checklist ("- [ ] TODO
        # write the runbook") into an "Optimist of the Year" score, which is the wrong signal.
        if os.path.splitext(f)[1].lower() not in (".md", ".txt"):
            todo_count += len(TODO_RE.findall(blob))
        # _emoji(), not EMOJI_RE directly -- the raw regex skips the typography denylist, so
        # this counter was inconsistent with the deck and commit-message counters it is
        # summed with.
        emoji_in_code += _emoji(blob)[0]

    deps_count = None
    nested_manifests = [f for f in tracked
                        if os.path.basename(f) in ("requirements.txt", "package.json", "pyproject.toml")
                        and "package-lock" not in f]
    counters = {
        "requirements.txt": lambda t: len([l for l in t.splitlines()
                                           if l.strip() and not l.startswith("#")]),
        "package.json": lambda t: len((json.loads(t).get("dependencies") or {}))
                                  + len((json.loads(t).get("devDependencies") or {})),
        "pyproject.toml": lambda t: t.lower().count("=") if "dependencies" in t.lower() else 0,
    }
    # Sum across EVERY manifest in scope -- a split backend/frontend project has two.
    for rel in nested_manifests:
        blob = _read_blob(repo, rel, ref_tip)
        if blob is None:
            continue
        try:
            deps_count = (deps_count or 0) + counters[os.path.basename(rel)](blob)
        except (ValueError, KeyError):
            pass

    tref = tracked if ref_tip else None
    claude_found = {k: _exists_any(repo, v, paths, tref) for k, v in CLAUDE_ARTIFACTS.items()}
    durability_found = {k: _exists_any(repo, v, paths, tref) for k, v in DURABILITY_ARTIFACTS.items()}

    # Root-only lookup misses nested layouts. Verified on a real submission: a project with
    # backend/requirements.txt + frontend/package.json reported NO dependency manifest,
    # which would have under-scored its Staying Power. Fall back to a basename match over
    # the tracked (already scope-filtered) file list.
    for key, candidates in DURABILITY_ARTIFACTS.items():
        if durability_found[key] is not None:
            continue
        wanted = {c.rsplit("/", 1)[-1].lower() for c in candidates}
        hits = [f for f in tracked if os.path.basename(f).lower() in wanted]
        if hits:
            durability_found[key] = {"path": hits[0], "is_dir": False, "entries": None,
                                     "scope": "submission" if paths else "repo-root",
                                     "found_nested": True, "all_paths": hits[:6]}
    # Nested CLAUDE.md / SKILL.md anywhere tracked (plugins, subprojects)
    nested_claude = [f for f in tracked if os.path.basename(f) in ("CLAUDE.md", "SKILL.md", "AGENTS.md")]

    total_int = int(total_all_time.strip()) if total_all_time and total_all_time.strip().isdigit() else None
    result.update(
        {
            "commits_in_window": len(commits),
            "commits_all_time": total_int,
            "commits_outside_window": (total_int - len(commits)) if total_int is not None else None,
            "merge_commits_in_window": int(merges.strip()) if merges and merges.strip().isdigit() else None,
            "authors_in_window": dict(sorted(authors.items(), key=lambda kv: -kv[1])),
            "distinct_authors": len(authors),
            "claude_trailer_commits": claude_commits,
            "first_commit_in_window": dates[0] if dates else None,
            "last_commit_in_window": dates[-1] if dates else None,
            "active_days_in_window": active_days,
            "insertions_in_window": ins,
            "deletions_in_window": dels,
            "files_touched_in_window": len(touched),
            "top_dirs_touched": dict(sorted(top_dirs.items(), key=lambda kv: -kv[1])[:12]),
            "tracked_file_count": len(tracked) if not tracked_err else None,
            "tracked_by_extension": dict(sorted(ext_counts.items(), key=lambda kv: -kv[1])[:12]),
            "test_files": test_files[:40],
            "test_file_count": len(test_files),
            "claude_code_artifacts": claude_found,
            "durability_artifacts": durability_found,
            "nested_agent_docs": nested_claude[:40],
            # --- superlative-award counters (fun awards only) ---
            "commit_hour_histogram": {str(h): n for h, n in hour_hist.items() if n},
            "nocturnal_commits_0000_0459": sum(hour_hist[h] for h in range(0, 5)),
            "commit_message_emoji": msg_emoji_total,
            "commit_message_emoji_tally": dict(sorted(msg_emoji_tally.items(), key=lambda kv: -kv[1])[:20]),
            "commits_with_emoji": emoji_commits,
            "emoji_in_code_or_docs": emoji_in_code,
            "revert_commits": reverts,
            "longest_commit_message": longest_msg,
            "biggest_single_commit": biggest,
            "tooling_path_touches": tooling_touches,
            "product_path_touches": product_touches,
            "yak_shave_ratio": round(tooling_touches / product_touches, 2) if product_touches else None,
            "todo_fixme_count": todo_count,
            "files_scanned_for_todo": scanned,
            "declared_dependency_count": deps_count,
            "commits": commits[:400],
        }
    )
    return result


# ------------------------------------------------- superlative roll-up ------------

def roll_up_superlatives(ev):
    """Flatten deck + all repos into one comparable scalar per superlative metric.

    Everything here is a mechanical counter for the FUN awards. None of it may inform a
    1-3 rubric band -- see the firewall rule in SKILL.md. `null` means "not measured",
    which is NOT the same as zero and must never be ranked as if it were.
    """
    d = ev.get("deck") or {}
    # A deck that FAILED to read is {"error": ...} -- truthy, but with no counters. Gate on
    # slide_count so a read failure yields null metrics rather than a confident-looking 0.
    deck_ok = bool(d.get("slide_count"))
    repos = [r for r in ev.get("repos", []) if not r.get("error")]
    sr = ev.get("self_reported", {})

    def s(key, default=0):
        return sum((r.get(key) or 0) for r in repos) if repos else default

    all_days = sorted({day for r in repos for day in (r.get("active_days_in_window") or [])})
    churn = s("insertions_in_window") + s("deletions_in_window")
    hours = {}
    for r in repos:
        for h, n in (r.get("commit_hour_histogram") or {}).items():
            hours[h] = hours.get(h, 0) + n
    biggest = max((r.get("biggest_single_commit") or {"churn": 0} for r in repos),
                  key=lambda b: b.get("churn") or 0, default={"churn": 0})
    longest_msg = max((r.get("longest_commit_message") or {"chars": 0} for r in repos),
                      key=lambda b: b.get("chars") or 0, default={"chars": 0})
    deck_emoji = d.get("emoji_count")
    msg_emoji = s("commit_message_emoji") if repos else None
    code_emoji = s("emoji_in_code_or_docs") if repos else None
    # Only the team's OWN harness artifacts count -- a root-scoped hit on a shared
    # monorepo is the organisers' file, not theirs.
    harness_hits = sum(1 for r in repos
                       for v in (r.get("claude_code_artifacts") or {}).values()
                       if v and (v.get("scope") == "submission" or not r.get("scoped_to")))
    deps = [r.get("declared_dependency_count") for r in repos
            if r.get("declared_dependency_count") is not None]

    return {
        # presentation shape
        "slides": d.get("slide_count"),
        "deck_words": d.get("total_words"),
        "densest_slide_words": d.get("densest_slide_words"),
        "talk_minutes": sr.get("talk_minutes") if sr.get("talk_minutes") is not None
                        else d.get("est_speaking_minutes"),
        "talk_minutes_is_observed": sr.get("talk_minutes") is not None,
        # register / vibe
        "emoji_total": (sum(x for x in (deck_emoji, msg_emoji, code_emoji) if x is not None)
                        if any(x is not None for x in (deck_emoji, msg_emoji, code_emoji)) else None),
        "emoji_deck": deck_emoji,
        "emoji_commits_and_code": (msg_emoji or 0) + (code_emoji or 0) if repos else None,
        "emoji_per_slide": d.get("emoji_per_slide"),
        "geek_score": d.get("geek_score"),
        "nerd_score": d.get("nerd_score"),
        "normie_score": d.get("normie_score"),
        "normalist_score": ((d.get("normie_score") or 0) - (d.get("geek_score") or 0)
                            - (d.get("nerd_score") or 0)) if deck_ok else None,
        "acronym_tokens": d.get("acronym_tokens"),
        # effort shape
        "tokens_reported": sr.get("tokens"),
        "commits": s("commits_in_window") if repos else None,
        "churn": churn if repos else None,
        "files_touched": s("files_touched_in_window") if repos else None,
        # Union, not max: two repos with one author each may be two different people, and
        # max() would report a one-person team.
        "distinct_authors": (len({a for r in repos for a in (r.get("authors_in_window") or {})})
                             if repos else None),
        "claude_trailered_commits": s("claude_trailer_commits") if repos else None,
        "active_days": len(all_days) if repos else None,
        "nocturnal_commits": s("nocturnal_commits_0000_0459") if repos else None,
        "peak_commit_hour": max(hours.items(), key=lambda kv: kv[1])[0] if hours else None,
        # craft residue
        "yak_shave_ratio": (round(s("tooling_path_touches") / s("product_path_touches"), 2)
                            if repos and s("product_path_touches") else None),
        "biggest_commit_churn": biggest.get("churn") if repos else None,
        "biggest_commit_share": (round((biggest.get("churn") or 0) / churn, 2)
                                 if repos and churn else None),
        "longest_commit_message_chars": longest_msg.get("chars") if repos else None,
        "revert_commits": s("revert_commits") if repos else None,
        "todo_fixme_count": s("todo_fixme_count") if repos else None,
        "declared_dependency_count": sum(deps) if deps else None,
        "harness_artifact_kinds": harness_hits if repos else None,
        "test_file_count": s("test_file_count") if repos else None,
    }


# ---------------------------------------------------------------- main ------------

def main(argv=None):
    ap = argparse.ArgumentParser(description="Collect hackathon judging evidence (no scoring).")
    ap.add_argument("--project", required=True, help="Submission name")
    ap.add_argument("--deck", default=None, help="Path to .pptx")
    ap.add_argument("--repo", action="append", default=[], help="Local repo path (repeatable)")
    ap.add_argument("--paths", action="append", default=[],
                    help="Pathspec to scope a shared monorepo to this submission, e.g. "
                         "reclass-intake (repeatable; applies to every --repo)")
    ap.add_argument("--since", default=None, help="Hackathon window start, e.g. 2026-08-11")
    ap.add_argument("--until", default=None, help="Hackathon window end")
    ap.add_argument("--team", default=None, help="Team name, if different from project")
    ap.add_argument("--ref", default=None,
                    help="Git ref to scope history to, e.g. origin/my-branch. Default: all "
                         "branches. REQUIRED when a team submitted on its own branch of a "
                         "shared repo, or its history blends with everyone else's.")
    ap.add_argument("--tokens", type=int, default=None,
                    help="Self-reported Claude token spend (for the Tokenmaxxer award)")
    ap.add_argument("--talk-minutes", type=float, default=None,
                    help="Observed presentation length in minutes (beats the word-count estimate)")
    ap.add_argument("--out", required=True, help="Output JSON path")
    args = ap.parse_args(argv)

    evidence = {
        "project": args.project,
        "team": args.team or args.project,
        "collected_at": datetime.now(timezone.utc).isoformat(),
        "collector_version": "1.1",
        "window": {"since": args.since, "until": args.until},
        "self_reported": {"tokens": args.tokens, "talk_minutes": args.talk_minutes},
        "deck": read_deck(args.deck),
        "repos": [read_repo(r, args.since, args.until, args.paths, args.ref) for r in args.repo],
        "missing_inputs": [],
    }
    evidence["superlative_metrics"] = roll_up_superlatives(evidence)
    if not args.deck:
        evidence["missing_inputs"].append("deck")
    if not args.repo:
        evidence["missing_inputs"].append("repo")
    if not args.since:
        evidence["missing_inputs"].append("window_start (cannot separate hackathon work from pre-existing work)")

    os.makedirs(os.path.dirname(os.path.abspath(args.out)) or ".", exist_ok=True)
    with open(args.out, "w", encoding="utf-8") as fh:
        json.dump(evidence, fh, indent=2, ensure_ascii=False)

    d = evidence["deck"] or {}
    print(f"project={args.project}")
    print(f"deck: {d.get('slide_count', 0)} slides, {len(d.get('quantified_claims', []))} quantified claims, "
          f"{len(d.get('unique_links', []))} links" + (f" ERROR={d.get('error')}" if d.get("error") else ""))
    for r in evidence["repos"]:
        if r.get("error"):
            print(f"repo {r['path']}: ERROR {r['error']}")
            continue
        print(f"repo {os.path.basename(r['path'])}: {r['commits_in_window']} commits in window "
              f"({r['commits_outside_window']} outside), {r['distinct_authors']} authors, "
              f"{r['claude_trailer_commits']} Claude-trailered, {r['test_file_count']} test files")
        if r.get("scope_warning"):
            print(f"  !! SCOPE: {r['scope_warning']}")
    m = evidence["superlative_metrics"]
    print(f"superlatives: slides={m['slides']} talk~{m['talk_minutes']}min "
          f"emoji={m['emoji_total']} geek={m['geek_score']} nerd={m['nerd_score']} "
          f"normalist={m['normalist_score']} nocturnal={m['nocturnal_commits']} "
          f"yak_shave={m['yak_shave_ratio']} tokens={m['tokens_reported']}")
    if evidence["missing_inputs"]:
        print("MISSING: " + "; ".join(evidence["missing_inputs"]))
    print(f"wrote {args.out}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
